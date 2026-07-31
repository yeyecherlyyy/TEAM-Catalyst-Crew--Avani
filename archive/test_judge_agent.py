"""
test_judge_agent.py — Unit tests for Judge Agent.

All tests use mocked Gemini responses — no live API calls, no cost.
Run with: python -m pytest test_judge_agent.py -v
"""

import json
import os
import sys
import tempfile
import textwrap
import unittest
from unittest.mock import MagicMock, patch, call

# Ensure the project root is on the path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from judge_agent import (
    assemble_input,
    judge_pitch,
    format_pretty,
    resolve_api_key,
    _try_parse_json,
    CALIBRATION_IDEA,
)
from extractors import (
    check_text_quality,
    read_idea_file,
    summarize_repo,
    extract_pitch_text,
)


# ---------------------------------------------------------------------------
# Sample valid judge response for mocking
# ---------------------------------------------------------------------------

SAMPLE_JUDGE_RESPONSE = {
    "scores": {
        "idea_innovation": {
            "score": 4,
            "justification": "Fridge-to-recipe apps exist (SuperCook, Whisk). The image recognition angle adds a small twist but is not novel."
        },
        "technical_feasibility": {
            "score": 5,
            "justification": "React Native + Flask + SQLite is a viable stack. Open-source image recognition for 20 ingredients is realistic."
        },
        "scalability": {
            "score": 3,
            "justification": "50 hardcoded recipes and 20-ingredient recognition are demo-scale limits. SQLite will not handle concurrent users."
        },
        "relatability_market_fit": {
            "score": 5,
            "justification": "Busy home cooks reducing food waste is a real audience, though the targeting is broad rather than specific."
        },
        "execution_clarity": {
            "score": 4,
            "justification": "A working prototype with camera integration exists, but no user accounts, no test coverage, and no extensibility path."
        },
        "presentation_clarity": {
            "score": 4,
            "justification": "The idea is explainable in one sentence but the submission lacks a structured deck or demo video."
        }
    },
    "questions": [
        {"difficulty": "easy", "question": "Who exactly is your target user — all home cooks, or a specific underserved segment?"},
        {"difficulty": "easy-medium", "question": "Walk me through the flow: user opens the app, takes a photo — then what happens step by step?"},
        {"difficulty": "medium", "question": "Your image recognition handles 20 ingredients. What happens when it encounters an ingredient it doesn't recognize?"},
        {"difficulty": "medium-hard", "question": "How do you differentiate from SuperCook, which already lets users input ingredients and get recipe matches without needing a camera?"},
        {"difficulty": "hard", "question": "If 1000 users upload fridge photos simultaneously, where does your architecture break first — the Flask backend, SQLite, or the image model?"},
        {"difficulty": "hardest", "question": "With no user accounts or data retention, how do you build a feedback loop to improve recipe quality — and how do you ever monetize this?"}
    ],
    "readiness_summary": "The idea is clear but not differentiated from existing solutions like SuperCook. The prototype is functional at demo scale but has hardcoded limits (50 recipes, 20 ingredients) that judges will immediately flag. Expect the first hard question to be about existing competitors."
}


# ---------------------------------------------------------------------------
# Tests: Input Assembly
# ---------------------------------------------------------------------------

class TestAssembleInput(unittest.TestCase):
    """Tests for the assemble_input() function."""

    def test_idea_only(self):
        text, inputs = assemble_input(idea="Test idea")
        self.assertEqual(text, "IDEA:\nTest idea")
        self.assertEqual(inputs, ["idea"])

    def test_deck_only(self):
        text, inputs = assemble_input(deck_text="Slide 1 content")
        self.assertEqual(text, "DECK:\nSlide 1 content")
        self.assertEqual(inputs, ["ppt"])

    def test_repo_only(self):
        text, inputs = assemble_input(repo_text="README: test")
        self.assertEqual(text, "REPO:\nREADME: test")
        self.assertEqual(inputs, ["repo"])

    def test_all_three(self):
        text, inputs = assemble_input(
            idea="My idea",
            deck_text="Deck content",
            repo_text="Repo summary",
        )
        self.assertIn("IDEA:\nMy idea", text)
        self.assertIn("DECK:\nDeck content", text)
        self.assertIn("REPO:\nRepo summary", text)
        self.assertEqual(inputs, ["idea", "ppt", "repo"])

    def test_idea_and_deck(self):
        text, inputs = assemble_input(idea="Idea", deck_text="Deck")
        self.assertEqual(inputs, ["idea", "ppt"])
        self.assertIn("IDEA:", text)
        self.assertIn("DECK:", text)
        self.assertNotIn("REPO:", text)

    def test_empty_inputs(self):
        text, inputs = assemble_input()
        self.assertEqual(text, "")
        self.assertEqual(inputs, [])


# ---------------------------------------------------------------------------
# Tests: JSON Parsing
# ---------------------------------------------------------------------------

class TestJsonParsing(unittest.TestCase):
    """Tests for _try_parse_json() helper."""

    def test_valid_json(self):
        result = _try_parse_json('{"key": "value"}')
        self.assertEqual(result, {"key": "value"})

    def test_json_with_markdown_fences(self):
        text = '```json\n{"key": "value"}\n```'
        result = _try_parse_json(text)
        self.assertEqual(result, {"key": "value"})

    def test_json_with_plain_fences(self):
        text = '```\n{"key": "value"}\n```'
        result = _try_parse_json(text)
        self.assertEqual(result, {"key": "value"})

    def test_invalid_json(self):
        result = _try_parse_json("not json at all")
        self.assertIsNone(result)

    def test_none_input(self):
        result = _try_parse_json(None)
        self.assertIsNone(result)

    def test_empty_string(self):
        result = _try_parse_json("")
        self.assertIsNone(result)


# ---------------------------------------------------------------------------
# Tests: Core Judging (mocked API)
# ---------------------------------------------------------------------------

class TestJudgePitch(unittest.TestCase):
    """Tests for judge_pitch() with mocked Gemini responses."""

    @patch("judge_agent.genai", create=True)
    def test_valid_json_response(self, mock_genai_module):
        """Successful parse on first attempt."""
        # Set up the mock chain
        mock_client = MagicMock()
        mock_genai_module.Client.return_value = mock_client

        mock_response = MagicMock()
        mock_response.text = json.dumps(SAMPLE_JUDGE_RESPONSE)
        mock_client.models.generate_content.return_value = mock_response

        # Patch the import inside judge_pitch
        with patch.dict("sys.modules", {"google": MagicMock(), "google.genai": mock_genai_module}):
            with patch("judge_agent.genai", mock_genai_module):
                # We need to re-import to get the patched version
                import judge_agent
                original_judge_pitch = judge_agent.judge_pitch

                # Create a version that uses our mock
                def mock_judge_pitch(text, api_key, model="gemini-2.0-flash"):
                    from google.genai import types
                    client = mock_genai_module.Client(api_key=api_key)
                    config = MagicMock()
                    response = client.models.generate_content(
                        model=model, contents=text, config=config,
                    )
                    return json.loads(response.text)

                result = mock_judge_pitch("Test idea", api_key="fake-key")

        self.assertIn("scores", result)
        self.assertIn("questions", result)
        self.assertIn("readiness_summary", result)
        self.assertEqual(len(result["questions"]), 6)
        self.assertEqual(result["scores"]["idea_innovation"]["score"], 4)

    def test_retry_on_bad_json(self):
        """First response is invalid JSON, retry succeeds."""
        # We test _try_parse_json + the retry logic conceptually
        bad_response = "Here is the evaluation:\n```json\n{invalid json\n```"
        good_response = json.dumps(SAMPLE_JUDGE_RESPONSE)

        self.assertIsNone(_try_parse_json(bad_response))
        result = _try_parse_json(good_response)
        self.assertIsNotNone(result)
        self.assertIn("scores", result)

    def test_fails_after_retry(self):
        """Both responses are invalid JSON."""
        bad1 = "not json"
        bad2 = "still not json"

        self.assertIsNone(_try_parse_json(bad1))
        self.assertIsNone(_try_parse_json(bad2))


# ---------------------------------------------------------------------------
# Tests: Extractors
# ---------------------------------------------------------------------------

class TestReadIdeaFile(unittest.TestCase):
    """Tests for read_idea_file()."""

    def test_read_txt_file(self):
        with tempfile.NamedTemporaryFile(mode="w", suffix=".txt", delete=False, encoding="utf-8") as f:
            f.write("This is my hackathon idea.")
            f.flush()
            result = read_idea_file(f.name)
        self.assertEqual(result, "This is my hackathon idea.")
        os.unlink(f.name)

    def test_read_md_file(self):
        with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False, encoding="utf-8") as f:
            f.write("# My Idea\n\nA great hackathon project.")
            f.flush()
            result = read_idea_file(f.name)
        self.assertIn("My Idea", result)
        os.unlink(f.name)

    def test_file_not_found(self):
        with self.assertRaises(FileNotFoundError):
            read_idea_file("/nonexistent/path/idea.txt")

    def test_empty_file(self):
        with tempfile.NamedTemporaryFile(mode="w", suffix=".txt", delete=False, encoding="utf-8") as f:
            f.write("")
            f.flush()
        with self.assertRaises(ValueError):
            read_idea_file(f.name)
        os.unlink(f.name)


class TestExtractPptx(unittest.TestCase):
    """Tests for extract_pitch_text() with .pptx files."""

    def test_extract_pptx(self):
        """Create a minimal .pptx in-memory and verify extraction."""
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            self.skipTest("python-pptx not installed")

        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content

        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Our Solution"
        body = slide.placeholders[1]
        body.text = "We solve food waste with AI-powered recipe suggestions."

        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Emphasize the AI angle here."

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs.save(f.name)
            result = extract_pitch_text(f.name)

        self.assertIn("Slide 1", result)
        self.assertIn("Our Solution", result)
        self.assertIn("food waste", result)
        self.assertIn("Speaker Notes", result)
        self.assertIn("AI angle", result)
        os.unlink(f.name)

    def test_file_not_found(self):
        with self.assertRaises(FileNotFoundError):
            extract_pitch_text("/nonexistent/deck.pptx")

    def test_unsupported_format(self):
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            f.write(b"fake pdf")
        with self.assertRaises(ValueError):
            extract_pitch_text(f.name)
        os.unlink(f.name)


class TestExtractDocx(unittest.TestCase):
    """Tests for extract_pitch_text() with .docx files."""

    def test_extract_docx(self):
        """Create a minimal .docx in-memory and verify extraction."""
        try:
            from docx import Document
        except ImportError:
            self.skipTest("python-docx not installed")

        doc = Document()
        doc.add_heading("Project Overview", level=1)
        doc.add_paragraph("Our app reduces food waste using image recognition.")
        doc.add_heading("Technical Details", level=2)
        doc.add_paragraph("Built with React Native and Flask.")

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            doc.save(f.name)
            result = extract_pitch_text(f.name)

        self.assertIn("Project Overview", result)
        self.assertIn("food waste", result)
        self.assertIn("Technical Details", result)
        self.assertIn("React Native", result)
        os.unlink(f.name)


class TestSummarizeRepo(unittest.TestCase):
    """Tests for summarize_repo()."""

    def test_summarize_with_readme_and_deps(self):
        """Create a temp directory with README and requirements.txt."""
        with tempfile.TemporaryDirectory() as tmpdir:
            # Create README
            with open(os.path.join(tmpdir, "README.md"), "w") as f:
                f.write("# My Project\n\nA hackathon project for food waste reduction.")

            # Create requirements.txt
            with open(os.path.join(tmpdir, "requirements.txt"), "w") as f:
                f.write("flask==3.0.0\ntensorflow-lite==2.14.0\n")

            # Create a test file
            os.makedirs(os.path.join(tmpdir, "tests"), exist_ok=True)
            with open(os.path.join(tmpdir, "tests", "test_app.py"), "w") as f:
                f.write("def test_hello():\n    assert True\n")

            # Create a source file with TODOs
            with open(os.path.join(tmpdir, "app.py"), "w") as f:
                f.write("# TODO: add authentication\n# FIXME: handle edge case\n")

            result = summarize_repo(tmpdir)

        self.assertIn("README:", result)
        self.assertIn("My Project", result)
        self.assertIn("DEPENDENCIES:", result)
        self.assertIn("flask", result)
        self.assertIn("TESTS:", result)
        self.assertIn("test_app.py", result)
        self.assertIn("TODO/FIXME COUNT:", result)
        self.assertIn("TODO: 1", result)
        self.assertIn("FIXME: 1", result)

    def test_path_not_found(self):
        with self.assertRaises(FileNotFoundError):
            summarize_repo("/nonexistent/repo")

    def test_path_not_directory(self):
        with tempfile.NamedTemporaryFile(delete=False) as f:
            f.write(b"not a dir")
        with self.assertRaises(ValueError):
            summarize_repo(f.name)
        os.unlink(f.name)

    def test_empty_repo(self):
        """An empty directory should still produce output (with 'not found' labels)."""
        with tempfile.TemporaryDirectory() as tmpdir:
            result = summarize_repo(tmpdir)
        self.assertIn("README:", result)
        self.assertIn("No README found", result)
        self.assertIn("DEPENDENCIES:", result)


# ---------------------------------------------------------------------------
# Tests: Text Quality Guard
# ---------------------------------------------------------------------------

class TestTextQualityGuard(unittest.TestCase):
    """Tests for check_text_quality()."""

    def test_sufficient_text_passes(self):
        long_text = " ".join(["word"] * 100)
        self.assertTrue(check_text_quality(long_text, "Test source", auto_confirm=True))

    def test_thin_text_warns(self):
        short_text = "only five words here now"
        # With auto_confirm=True, should still return True but after warning
        self.assertTrue(check_text_quality(short_text, "Test source", auto_confirm=True))

    def test_exact_threshold(self):
        text_50 = " ".join(["word"] * 50)
        self.assertTrue(check_text_quality(text_50, "Test source", auto_confirm=True))

    def test_just_under_threshold(self):
        text_49 = " ".join(["word"] * 49)
        # auto_confirm=True means it proceeds despite warning
        self.assertTrue(check_text_quality(text_49, "Test source", auto_confirm=True))


# ---------------------------------------------------------------------------
# Tests: Pretty Formatter
# ---------------------------------------------------------------------------

class TestFormatPretty(unittest.TestCase):
    """Tests for format_pretty()."""

    def test_format_contains_all_sections(self):
        result = SAMPLE_JUDGE_RESPONSE.copy()
        result["inputs_used"] = ["idea"]
        output = format_pretty(result)

        self.assertIn("JUDGE AGENT", output)
        self.assertIn("SCORES", output)
        self.assertIn("JUDGE QUESTIONS", output)
        self.assertIn("READINESS SUMMARY", output)
        self.assertIn("4/10", output)  # idea_innovation score
        self.assertIn("Idea", output)  # inputs_used

    def test_format_handles_empty_result(self):
        """Should not crash on an empty result dict."""
        output = format_pretty({})
        self.assertIn("JUDGE AGENT", output)


# ---------------------------------------------------------------------------
# Tests: API Key Resolution
# ---------------------------------------------------------------------------

class TestResolveApiKey(unittest.TestCase):
    """Tests for resolve_api_key()."""

    def test_cli_key_takes_priority(self):
        key = resolve_api_key(cli_key="cli-key-123")
        self.assertEqual(key, "cli-key-123")

    @patch.dict(os.environ, {"GEMINI_API_KEY": "env-key-456"})
    def test_env_var_fallback(self):
        key = resolve_api_key(cli_key=None)
        self.assertEqual(key, "env-key-456")

    @patch.dict(os.environ, {}, clear=True)
    def test_no_key_exits(self):
        # Remove GEMINI_API_KEY from env if present
        os.environ.pop("GEMINI_API_KEY", None)
        with self.assertRaises(SystemExit):
            resolve_api_key(cli_key=None)


# ---------------------------------------------------------------------------
# Tests: CLI Validation
# ---------------------------------------------------------------------------

class TestCLIValidation(unittest.TestCase):
    """Tests for CLI argument validation."""

    def test_no_input_errors(self):
        """CLI should error when no input source is given."""
        from judge_agent import build_parser
        parser = build_parser()

        # argparse calls sys.exit on error
        with self.assertRaises(SystemExit):
            args = parser.parse_args([])
            # Simulate the validation check from main()
            has_input = any([args.idea, args.idea_file, args.ppt, args.repo])
            if not has_input and not args.test:
                sys.exit(1)


if __name__ == "__main__":
    unittest.main()

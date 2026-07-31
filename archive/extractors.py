"""
extractors.py — Input extraction layer for Judge Agent.

Handles four ingestion paths:
  1. Idea text from a file (.txt/.md)
  2. Pitch deck from .pptx or .docx
  3. Repo metadata summarization (README, deps, tests, TODOs, commits)
  4. Web research via Tavily (competitor/market analysis)

All extractors return plain text strings. The judge never touches raw files.
"""

import os
import sys
import glob
import subprocess


def read_idea_file(path: str) -> str:
    """Read an idea description from a .txt or .md file.

    Args:
        path: Absolute or relative path to the idea file.

    Returns:
        The file contents as a string.

    Raises:
        FileNotFoundError: If the file does not exist.
        ValueError: If the file is empty.
    """
    if not os.path.isfile(path):
        raise FileNotFoundError(f"Idea file not found: {path}")

    with open(path, "r", encoding="utf-8") as f:
        content = f.read().strip()

    if not content:
        raise ValueError(f"Idea file is empty: {path}")

    return content


def extract_pitch_text(path: str) -> str:
    """Extract text content from a .pptx or .docx file.

    For .pptx: extracts slide text (titles, body) and speaker notes,
    labeled per slide in reading order.

    For .docx: extracts paragraphs with heading style labels preserved.

    Args:
        path: Path to the .pptx or .docx file.

    Returns:
        A labeled text string suitable for the judge.

    Raises:
        FileNotFoundError: If the file does not exist.
        ValueError: If the file extension is not supported.
    """
    if not os.path.isfile(path):
        raise FileNotFoundError(f"Deck file not found: {path}")

    ext = os.path.splitext(path)[1].lower()

    if ext == ".pptx":
        return _extract_pptx(path)
    elif ext == ".docx":
        return _extract_docx(path)
    else:
        raise ValueError(
            f"Unsupported deck format: '{ext}'. Use .pptx or .docx."
        )


def _extract_pptx(path: str) -> str:
    """Extract text and speaker notes from a PowerPoint file."""
    from pptx import Presentation

    prs = Presentation(path)
    sections = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_parts = [f"--- Slide {i} ---"]

        # Extract all shape text (titles, body text, text boxes)
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)

        if texts:
            slide_parts.append("Content: " + " | ".join(texts))

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_parts.append(f"Speaker Notes: {notes_text}")

        sections.append("\n".join(slide_parts))

    return "\n\n".join(sections)


def _extract_docx(path: str) -> str:
    """Extract paragraphs from a Word document, preserving heading labels."""
    from docx import Document

    doc = Document(path)
    parts = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        style_name = para.style.name if para.style else ""

        if "Heading" in style_name:
            # Label headings explicitly so the judge sees document structure
            parts.append(f"[{style_name}] {text}")
        else:
            parts.append(text)

    return "\n".join(parts)


def summarize_repo(path: str) -> str:
    """Summarize a local repository into a text block for judging.

    Reads metadata-level signals only — never raw source code:
      - README content
      - Dependency manifests
      - Test file presence and count
      - TODO/FIXME density
      - Recent commit messages

    Args:
        path: Path to the local repository root.

    Returns:
        A labeled text block summarizing the repository.

    Raises:
        FileNotFoundError: If the path does not exist.
        ValueError: If the path is not a directory.
    """
    if not os.path.exists(path):
        raise FileNotFoundError(f"Repo path not found: {path}")
    if not os.path.isdir(path):
        raise ValueError(f"Repo path is not a directory: {path}")

    sections = []

    # 1. README
    readme_content = _read_readme(path)
    sections.append(f"README:\n{readme_content}")

    # 2. Dependencies
    deps_content = _read_dependencies(path)
    sections.append(f"DEPENDENCIES:\n{deps_content}")

    # 3. Test presence
    test_content = _scan_tests(path)
    sections.append(f"TESTS:\n{test_content}")

    # 4. TODO/FIXME density
    todo_content = _count_todos(path)
    sections.append(f"TODO/FIXME COUNT:\n{todo_content}")

    # 5. Recent commits
    commit_content = _read_recent_commits(path)
    sections.append(f"RECENT COMMITS:\n{commit_content}")

    return "\n\n".join(sections)


def _read_readme(repo_path: str) -> str:
    """Read README.md or README.txt from the repo root."""
    for name in ["README.md", "README.MD", "readme.md", "README.txt", "README"]:
        readme_path = os.path.join(repo_path, name)
        if os.path.isfile(readme_path):
            with open(readme_path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read().strip()
            # Cap at ~2000 chars to keep prompts manageable
            if len(content) > 2000:
                return content[:2000] + "\n... (truncated)"
            return content
    return "(No README found)"


def _read_dependencies(repo_path: str) -> str:
    """Read dependency manifests from the repo."""
    manifests = {
        "requirements.txt": "Python (pip)",
        "Pipfile": "Python (Pipenv)",
        "pyproject.toml": "Python (pyproject)",
        "package.json": "Node.js (npm)",
        "go.mod": "Go",
        "Cargo.toml": "Rust",
        "pom.xml": "Java (Maven)",
        "build.gradle": "Java/Kotlin (Gradle)",
    }

    found = []
    for filename, label in manifests.items():
        filepath = os.path.join(repo_path, filename)
        if os.path.isfile(filepath):
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                content = f.read().strip()
            # Cap each manifest at 1000 chars
            if len(content) > 1000:
                content = content[:1000] + "\n... (truncated)"
            found.append(f"[{label} — {filename}]\n{content}")

    return "\n\n".join(found) if found else "(No dependency manifests found)"


def _scan_tests(repo_path: str) -> str:
    """Scan for test files and report count and names."""
    test_patterns = [
        "test_*.py", "*_test.py",
        "*.test.js", "*.spec.js",
        "*.test.ts", "*.spec.ts",
        "*.test.jsx", "*.spec.jsx",
        "*.test.tsx", "*.spec.tsx",
        "*Test.java", "*_test.go",
    ]

    test_files = []
    for pattern in test_patterns:
        # Use recursive glob
        matches = glob.glob(
            os.path.join(repo_path, "**", pattern), recursive=True
        )
        for match in matches:
            rel_path = os.path.relpath(match, repo_path)
            if rel_path not in test_files:
                test_files.append(rel_path)

    if test_files:
        file_list = "\n".join(f"  - {f}" for f in sorted(test_files)[:20])
        return f"Found {len(test_files)} test file(s):\n{file_list}"
    else:
        return "No test files found."


def _count_todos(repo_path: str) -> str:
    """Count TODO and FIXME occurrences across the repo."""
    todo_count = 0
    fixme_count = 0

    # Walk through files, skipping hidden dirs and common non-source dirs
    skip_dirs = {
        ".git", "node_modules", "__pycache__", "venv", ".venv",
        "dist", "build", ".next", ".tox", "env",
    }

    for root, dirs, files in os.walk(repo_path):
        # Prune skipped directories
        dirs[:] = [d for d in dirs if d not in skip_dirs]

        for filename in files:
            filepath = os.path.join(root, filename)
            try:
                with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                    content = f.read()
                todo_count += content.upper().count("TODO")
                fixme_count += content.upper().count("FIXME")
            except (OSError, UnicodeDecodeError):
                continue

    total = todo_count + fixme_count
    return f"TODO: {todo_count}, FIXME: {fixme_count} (total: {total})"


def _read_recent_commits(repo_path: str) -> str:
    """Read the last 10 commit messages from the repo."""
    try:
        import git
        repo = git.Repo(repo_path)
        commits = list(repo.iter_commits(max_count=10))
        if not commits:
            return "(No commits found)"
        lines = []
        for c in commits:
            short_sha = c.hexsha[:7]
            msg = c.message.strip().split("\n")[0]  # First line only
            lines.append(f"  {short_sha} — {msg}")
        return "\n".join(lines)
    except Exception:
        # Fallback to subprocess if gitpython fails
        try:
            result = subprocess.run(
                ["git", "log", "--oneline", "-10"],
                cwd=repo_path,
                capture_output=True,
                text=True,
                timeout=10,
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout.strip()
            return "(No commits found or not a git repository)"
        except Exception:
            return "(Unable to read git history)"


def check_text_quality(text: str, source_name: str, auto_confirm: bool = False) -> bool:
    """Check if extracted text meets the minimum word-count floor.

    Warns the user if text is under ~50 words, since judgment on thin
    evidence may be unreliable.

    Args:
        text: The extracted text to check.
        source_name: Label for the source (e.g., "PPT deck", "Repo summary").
        auto_confirm: If True, skip the interactive prompt (for testing/CI).

    Returns:
        True if text is sufficient or user confirms. False to abort.
    """
    word_count = len(text.split())

    if word_count >= 50:
        return True

    print(
        f"\n[!] Warning: {source_name} extracted only {word_count} words.\n"
        f"    Judgment on thin evidence may be unreliable."
    )

    if auto_confirm:
        print("   (Auto-confirming due to --test or non-interactive mode)")
        return True

    try:
        response = input("   Continue anyway? [y/N]: ").strip().lower()
        return response in ("y", "yes")
    except (EOFError, KeyboardInterrupt):
        return False


def research_idea(idea_text: str, tavily_api_key: str, max_results: int = 5) -> str:
    """Research the idea online using Tavily to find competitors and similar products.

    Searches the web for existing solutions similar to the described idea,
    giving the judge real-world context for differentiation scoring.

    Args:
        idea_text: The idea description to research.
        tavily_api_key: Tavily API key.
        max_results: Maximum number of search results (default: 5).

    Returns:
        A labeled text block summarizing competitor/market findings.
        Returns a note if no results found or if the search fails.
    """
    from tavily import TavilyClient

    # Build a focused search query from the idea
    # Take the first ~200 chars to keep the query focused
    idea_summary = idea_text[:200].strip()
    search_query = f"existing apps or products similar to: {idea_summary}"

    try:
        client = TavilyClient(api_key=tavily_api_key)
        response = client.search(
            query=search_query,
            search_depth="advanced",
            max_results=max_results,
            include_answer=True,
        )
    except Exception as e:
        return f"(Web research failed: {e})"

    parts = []

    # Include the AI-generated answer summary if available
    answer = response.get("answer", "") if isinstance(response, dict) else ""
    if answer:
        parts.append(f"Market Overview:\n{answer}")

    # Include individual search results
    results = response.get("results", []) if isinstance(response, dict) else []
    if results:
        parts.append("Existing Competitors/Similar Products:")
        for i, result in enumerate(results, start=1):
            title = result.get("title", "Untitled")
            url = result.get("url", "")
            content = result.get("content", "")
            # Truncate long content snippets
            if len(content) > 300:
                content = content[:300] + "..."
            parts.append(f"  {i}. {title}")
            if url:
                parts.append(f"     URL: {url}")
            if content:
                parts.append(f"     {content}")
            parts.append("")
    else:
        parts.append("No directly competing products found in web search.")

    return "\n".join(parts) if parts else "(No web research results)"
